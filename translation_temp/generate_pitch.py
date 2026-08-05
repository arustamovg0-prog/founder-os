import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_slide(prs, title, content, layout_idx=1):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    if layout_idx == 1:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        if isinstance(content, list):
            for i, point in enumerate(content):
                if i == 0:
                    tf.text = point
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
        else:
            tf.text = content
            
    return slide

def main():
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Founder OS / UNTITLED"
    subtitle.text = "Комплексная венчурная B2G/B2B-инфраструктура"
    
    # 2. Решаемая проблема
    create_slide(prs, "Решаемая проблема (Problem Solved)", [
        "Низкая готовность стартапов к инвестициям (Investment Readiness): Большинство проектов на ранних стадиях имеют слабые бизнес-модели, непросчитанную юнит-экономику и отсутствие упаковки, из-за чего фонды теряют до 90% потенциального сделочного потока (deal flow).",
        "Хаос и высокие затраты фондов на первичный фильтр: Инвесторские организации (Uz Combinator, Yoshlar Ventures и др.) сталкиваются с необработанным потоком сырых заявок, вручную тратя сотни часов на первичный скрининг и аудит.",
        "Отсутствие единой методологии отбора: На рынке нет прозрачного, пошагового инструмента, который бы объективно оценивал стартапы на каждом этапе их развития и защищал инвестора от риска вложиться в «дутый» проект."
    ])
    
    # 3. Краткое описание продукта
    create_slide(prs, "О продукте (Product Description)", [
        "Founder OS / UNTITLED — это комплексная венчурная B2G/B2B-инфраструктура.",
        "Она объединяет аналитическую систему скоринга стартапов и геймифицированный акселерационный конвейер для вывода проектов в состояние Investment Ready."
    ])

    # 4. Как это работает: Геймификация
    create_slide(prs, "Как это работает: Геймифицированный пайплайн", [
        "Level-up система в Founder OS: Весь путь подготовки выстроен по принципу видеоигры, где главная цель и финальный «Босс» — это инвестор.",
        "Независимо от того, насколько готовым считает себя стартап на входе, он начинает путь с фундамента (первого уровня) и поэтапно проходит всю нашу систему."
    ])

    # 5. Как это работает: Двойной фильтр
    create_slide(prs, "Как это работает: Двойной фильтр (AI + Эксперты)", [
        "Переход на каждый новый этап (level-up) не происходит автоматически.",
        "Шаг 1: Оценка AI-скоринга (UNTITLED).",
        "Шаг 2: Метрики, документы и результаты ИИ-оценки разбираются нашими экспертами.",
        "Только после успешной экспертной защиты стартапу открывается доступ к следующему уровню."
    ])

    # 6. Как это работает: Контроль качества
    create_slide(prs, "Как это работает: Контроль качества и доступ к капиталу", [
        "Этапность продолжается до тех пор, пока проект не будет полностью отфильтрован и «упакован».",
        "Только когда все уровни успешно пройдены и мы на 100% удостоверимся, что стартап готов к масштабированию, система «открывает двери».",
        "Фаундер получает прямой доступ к инвесторам и инвестиционным комитетам фондов."
    ])

    # 7. Product: For Founders
    create_slide(prs, "Платформа Founder OS: Для Фаундеров", [
        "Step-by-step Roadmap: Пошаговый трекинг развития стартапа.",
        "AI Copilot: Умный помощник для ответа на вопросы инвесторов и валидации гипотез.",
        "Centralized Data Room: Единое хранилище артефактов (Pitch Deck, Financial Model).",
        "Pitches: Управление запросами к инвесторам."
    ])

    # 8. Product: For Investors
    create_slide(prs, "Платформа Founder OS: Для Инвесторов", [
        "Curated Deal Flow: Доступ к на 100% подготовленным проектам.",
        "AI-Scoring: Мгновенный доступ к результатам независимого аудита.",
        "CRM Pipeline: Канбан-доска для управления переговорами.",
        "Portfolio Analytics: Дашборд с показателями проинвестированных компаний."
    ])

    # 9. Product: For Admins
    create_slide(prs, "Платформа Founder OS: Для Администраторов", [
        "Ecosystem Health: Интерактивный радар развития всех проектов.",
        "Stage Review: Центр верификации стартапов и экспертной защиты.",
        "Resident Database: Полный реестр всех стартапов-резидентов."
    ])

    # 10. Technology & Status
    create_slide(prs, "Технологии и Статус проекта", [
        "Modern Stack: Next.js 16 (App Router), React 19, Tailwind CSS 4, Firebase.",
        "AI Интеграция: Google Gemini API для генерации Executive Summary, скоринга и AI Hints.",
        "Current Status: Проект готов, развернут в Production (Vercel) и готов к пилотному запуску."
    ])

    # 11. Call to Action
    create_slide(prs, "Вместе строим инфраструктуру рынка!", [
        "Присоединяйтесь к Founder OS / UNTITLED.",
        "Почта: [Укажите email]",
        "Сайт: [Укажите сайт]"
    ])

    output_path = "Founder_OS_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    main()
