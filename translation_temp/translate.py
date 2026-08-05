import sys
import time
from pptx import Presentation
from deep_translator import GoogleTranslator

def translate_presentation(input_path, output_path):
    translator = GoogleTranslator(source='auto', target='en')
    prs = Presentation(input_path)
    
    translation_cache = {}

    def get_translation(text):
        text_strip = text.strip()
        if not text_strip:
            return text
        if not any(c.isalpha() for c in text_strip):
            return text
            
        if text_strip in translation_cache:
            return translation_cache[text_strip]
        try:
            translated = translator.translate(text_strip)
            translation_cache[text_strip] = translated
            return translated
        except Exception as e:
            print(f"Error translating '{text_strip}': {e}")
            return text

    def translate_shape(shape):
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text
                if text and any(c.isalpha() for c in text):
                    translated_text = get_translation(text)
                    if paragraph.runs:
                        paragraph.runs[0].text = translated_text
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""
        
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame") and cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            text = paragraph.text
                            if text and any(c.isalpha() for c in text):
                                translated_text = get_translation(text)
                                if paragraph.runs:
                                    paragraph.runs[0].text = translated_text
                                    for i in range(1, len(paragraph.runs)):
                                        paragraph.runs[i].text = ""

        if hasattr(shape, "shapes"): # GroupShape
            for child_shape in shape.shapes:
                translate_shape(child_shape)

    total_slides = len(prs.slides)
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Translating slide {slide_idx + 1}/{total_slides}...")
        for shape in slide.shapes:
            translate_shape(shape)

    prs.save(output_path)
    print("Translation completed successfully!")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python translate.py <input.pptx> <output.pptx>")
        sys.exit(1)
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    translate_presentation(input_file, output_file)
